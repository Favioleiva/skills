#!/usr/bin/env python3
"""
validate_pptx.py - Automated Quality Assurance auditor for assembled slide decks.
Part of the english-class-slides (v1.1) skill package.
"""

import os
import sys
import json
import argparse

if sys.stdout.encoding != 'utf-8':
    try:
        sys.stdout.reconfigure(encoding='utf-8')
    except Exception:
        pass

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("Error: python-pptx is not installed. Please run: pip install python-pptx")
    sys.exit(1)

def validate_presentation(lang_label, filepath, expected_count=53):
    print(f"\nAUDITING [{lang_label.upper()}]: {filepath}")
    if not os.path.exists(filepath):
        print(f"  [X] FAIL: Presentation file not found: {filepath}")
        return False

    prs = Presentation(filepath)
    slide_count = len(prs.slides)
    if slide_count != expected_count:
        print(f"  [X] FAIL: Expected {expected_count} slides, found {slide_count}")
        return False
    print(f"  [OK] Slide count: {slide_count}/{expected_count}")

    # Check Aspect Ratio (16:9)
    w_in = prs.slide_width / Inches(1)
    h_in = prs.slide_height / Inches(1)
    aspect_ratio = w_in / h_in
    if abs(aspect_ratio - (16/9)) > 0.02:
        print(f"  [X] FAIL: Non-16:9 aspect ratio ({aspect_ratio:.3f})")
        return False
    print(f"  [OK] Dimensions: 16:9 widescreen ({w_in:.2f}\" x {h_in:.2f}\")")

    failures = []
    below_24pt_count = 0

    for i, slide in enumerate(prs.slides, 1):
        shapes = list(slide.shapes)
        pics = [s for s in shapes if s.shape_type == 13]
        tbs = [s for s in shapes if s.has_text_frame]

        if not pics:
            failures.append(f"Slide {i}: Missing full-bleed background picture")
        if not tbs:
            failures.append(f"Slide {i}: Missing text box")
            continue

        if i == 1:
            # Cover Slide
            if len(tbs) < 2:
                failures.append(f"Slide 1 (Cover): Expected 2 text frames (Title + Credits), found {len(tbs)}")
        elif 2 <= i <= 51:
            # Content Slide
            p = tbs[0].text_frame.paragraphs[0]
            full_text = tbs[0].text_frame.text.strip()
            lines = [l for l in full_text.split("\n") if l.strip()]
            if len(lines) != 1:
                failures.append(f"Slide {i}: Expected exactly 1 sentence with zero title, found {len(lines)} lines")
            if p.font and p.font.size and p.font.size.pt < 24:
                below_24pt_count += 1
        elif 52 <= i <= 53:
            # Vocabulary Slide
            if len(tbs) < 2:
                failures.append(f"Slide {i} (Vocab): Expected 2 text frames (Heading + List), found {len(tbs)}")
            else:
                items = [p.text.strip() for p in tbs[1].text_frame.paragraphs if p.text.strip()]
                if len(items) != 10:
                    failures.append(f"Slide {i} (Vocab): Expected exactly 10 items, found {len(items)}")

    if failures:
        print(f"  [X] FAILURES DETECTED ({len(failures)}):")
        for f in failures:
            print(f"    - {f}")
        return False

    print(f"  [OK] Content slides with font < 24 pt: {below_24pt_count}")
    print(f"  [OK] ALL CHECKS PASSED FOR {lang_label.upper()}!")
    return True

def validate_kana_reading(filepath, expected_count=53):
    print(f"\nAUDITING [JAPANESE KANA READING LAYER]: {filepath}")
    if not os.path.exists(filepath):
        print(f"  [-] Notice: Kana reading file not found (optional): {filepath}")
        return True

    import re
    with open(filepath, "r", encoding="utf-8") as f:
        content = f.read()

    slides = re.split(r"(?:スライド)\s+(\d+)", content)
    slide_entries = []
    for i in range(1, len(slides), 2):
        s_num = int(slides[i])
        s_text = slides[i+1].strip()
        slide_entries.append((s_num, s_text))

    if len(slide_entries) != expected_count:
        print(f"  [X] FAIL: Expected {expected_count} slides, found {len(slide_entries)}")
        return False
    print(f"  [OK] Slide count: {len(slide_entries)}/{expected_count}")

    kanji_pattern = re.compile(r"[\u4e00-\u9fff\u3400-\u4dbf]")
    kanji_errors = []
    romaji_errors = []
    ruby_errors = []

    for s_num, s_text in slide_entries:
        for line in s_text.split("\n"):
            val = line.strip()
            if not val:
                continue
            for prefix in ["タイトル:", "サブタイトル:", "著者:", "AI支援制作:", "クレジット:", "本文:", "単語:"]:
                if val.startswith(prefix):
                    val = val[len(prefix):].strip()
                    break
            if re.match(r"^\d+\.\s*", val):
                val = re.sub(r"^\d+\.\s*", "", val)

            km = kanji_pattern.findall(val)
            if km:
                kanji_errors.append((s_num, val, km))
            rm = re.findall(r"[a-zA-Z]", val)
            if rm:
                romaji_errors.append((s_num, val, rm))
            rbm = re.findall(r"[\(（\[［\{｛][\u3040-\u309f\u30a0-\u30ff]+[\)）\]］\}｝]", val)
            if rbm:
                ruby_errors.append((s_num, val, rbm))

    if kanji_errors:
        print(f"  [X] FAIL: Unexpected Kanji detected in reading layer ({len(kanji_errors)} occurrences)")
        return False
    print(f"  [OK] Zero unexpected Kanji/CJK ideographs in phonetic text")

    if romaji_errors:
        print(f"  [X] FAIL: Unexpected Romaji detected in reading layer ({len(romaji_errors)} occurrences)")
        return False
    print(f"  [OK] Zero Romaji in phonetic text")

    if ruby_errors:
        print(f"  [X] FAIL: Ruby/Furigana annotation tags detected ({len(ruby_errors)} occurrences)")
        return False
    print(f"  [OK] Zero Ruby/annotation tags in phonetic text")
    print(f"  [OK] ALL CHECKS PASSED FOR JAPANESE KANA READING LAYER!")
    return True

def main():
    parser = argparse.ArgumentParser(description="Audit and validate generated PPTX presentations and kana reading layer.")
    parser.add_argument("--config", default="templates/deck_config.template.json", help="Path to deck config JSON")
    args = parser.parse_args()

    with open(args.config, "r", encoding="utf-8") as f:
        config = json.load(f)

    paths = config["paths"]
    expected = config.get("slide_count", 53)

    all_passed = True
    for lang, path_key in [("English", "output_pptx_en"), ("Spanish", "output_pptx_es"), ("Japanese", "output_pptx_ja")]:
        pptx_path = paths[path_key]
        ok = validate_presentation(lang, pptx_path, expected)
        if not ok:
            all_passed = False

    if "reading_ja_kana" in paths:
        ok_kana = validate_kana_reading(paths["reading_ja_kana"], expected)
        if not ok_kana:
            all_passed = False

    print("\n" + "=" * 60)
    if all_passed:
        print("[ALL PASS] FINAL VERIFICATION: ALL PRESENTATIONS AND KANA LAYER ARE 100% COMPLIANT!")
    else:
        print("[FAIL] VERIFICATION FAILED. Please review issues above.")
    print("=" * 60)

if __name__ == "__main__":
    main()
