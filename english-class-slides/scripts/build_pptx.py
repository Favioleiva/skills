#!/usr/bin/env python3
"""
build_pptx.py - Automated Google-Slides-compatible PowerPoint assembly pipeline.
Part of the english-class-slides (v1.1) skill package.
"""

import os
import sys
import re
import json
import argparse

if sys.stdout.encoding != 'utf-8':
    try:
        sys.stdout.reconfigure(encoding='utf-8')
    except Exception:
        pass

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx is not installed. Please run: pip install python-pptx")
    sys.exit(1)

# Default Slide dimensions: 16:9 widescreen
SLIDE_WIDTH = Inches(13.333333)
SLIDE_HEIGHT = Inches(7.5)

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def parse_slides_file(filepath):
    with open(filepath, "r", encoding="utf-8") as f:
        content = f.read()

    slides = []
    pattern = r"(?:SLIDE|DIAPOSITIVA|スライド)\s+(\d+)"
    splits = re.split(pattern, content)
    
    for i in range(1, len(splits), 2):
        slide_num = int(splits[i])
        slide_text = splits[i+1].strip()
        slides.append((slide_num, slide_text))
    
    return slides

def parse_cover_slide(slide_text, lang):
    lines = [l.strip() for l in slide_text.split("\n") if l.strip()]
    data = {
        "title": "",
        "subtitle": "",
        "author": "",
        "ai_production": "",
        "credits": []
    }
    
    mode = None
    for line in lines:
        if re.match(r"^(?:Design|Diseño|デザイン):", line, re.IGNORECASE):
            continue
        elif re.match(r"^(?:Title|Título|タイトル):", line, re.IGNORECASE):
            data["title"] = re.sub(r"^(?:Title|Título|タイトル):\s*", "", line, flags=re.IGNORECASE).strip()
            mode = None
        elif re.match(r"^(?:Subtitle|Subtítulo|サブタイトル):", line, re.IGNORECASE):
            data["subtitle"] = re.sub(r"^(?:Subtitle|Subtítulo|サブタイトル):\s*", "", line, flags=re.IGNORECASE).strip()
            mode = None
        elif re.match(r"^(?:Author|Autor|著者):", line, re.IGNORECASE):
            data["author"] = re.sub(r"^(?:Author|Autor|著者):\s*", "", line, flags=re.IGNORECASE).strip()
            mode = None
        elif re.match(r"^(?:AI-Assisted Production|Producción asistida por IA|AI支援制作):", line, re.IGNORECASE):
            data["ai_production"] = re.sub(r"^(?:AI-Assisted Production|Producción asistida por IA|AI支援制作):\s*", "", line, flags=re.IGNORECASE).strip()
            mode = None
        elif re.match(r"^(?:Credits|Créditos|クレジット):", line, re.IGNORECASE):
            mode = "credits"
        elif mode == "credits":
            data["credits"].append(line)
            
    return data

def parse_content_slide(slide_text):
    lines = [l.strip() for l in slide_text.split("\n") if l.strip()]
    sentence = ""
    for line in lines:
        if re.match(r"^(?:Design|Diseño|デザイン):", line, re.IGNORECASE):
            continue
        m = re.match(r"^(?:Text|Texto|本文):\s*(.*)", line, re.IGNORECASE)
        if m:
            sentence = m.group(1).strip()
            break
        elif not line.startswith("SLIDE") and not line.startswith("DIAPOSITIVA") and not line.startswith("スライド"):
            sentence = line
            break
    return sentence

def parse_vocab_slide(slide_text):
    lines = [l.strip() for l in slide_text.split("\n") if l.strip()]
    title = ""
    words = []
    mode = None
    for line in lines:
        if re.match(r"^(?:Design|Diseño|デザイン):", line, re.IGNORECASE):
            continue
        elif re.match(r"^(?:Title|Título|タイトル):", line, re.IGNORECASE):
            title = re.sub(r"^(?:Title|Título|タイトル):\s*", "", line, flags=re.IGNORECASE).strip()
        elif re.match(r"^(?:Words|Palabras|単語):", line, re.IGNORECASE):
            mode = "words"
        elif mode == "words" or re.match(r"^\d+\.", line):
            if re.match(r"^\d+\.", line):
                words.append(line)
    return title, words

def resolve_bg_path(slide_num, mapping, images_dir):
    for k, v in mapping.items():
        if "-" in str(k):
            start, end = map(int, str(k).split("-"))
            if start <= slide_num <= end:
                return os.path.join(images_dir, v)
        elif int(k) == slide_num:
            return os.path.join(images_dir, v)
    # Default fallback
    return os.path.join(images_dir, f"design_{slide_num:02d}.png")

def get_layout_override(slide_num, overrides):
    for k, v in overrides.items():
        if "-" in str(k):
            start, end = map(int, str(k).split("-"))
            if start <= slide_num <= end:
                return v
        elif int(k) == slide_num:
            return v
    return None

def create_deck(lang_code, txt_file, output_pptx, config):
    print(f"Building [{lang_code.upper()}] Presentation: {output_pptx}")
    slides_data = parse_slides_file(txt_file)
    expected_slides = config.get("slide_count", 53)
    if len(slides_data) != expected_slides:
        print(f"Warning: parsed {len(slides_data)} slides, expected {expected_slides}.")
    
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # Fonts
    if lang_code in ["en", "es"]:
        FONT_TITLE = "Cormorant Garamond"
        FONT_BODY = "Nunito Sans"
        FONT_AUTHOR = "Nunito Sans"
    else: # ja
        FONT_TITLE = "Noto Serif JP"
        FONT_BODY = "Noto Sans JP"
        FONT_AUTHOR = "Noto Sans JP"

    COLOR_TITLE = hex_to_rgb("102238")
    COLOR_BODY = hex_to_rgb("1A2433")
    COLOR_SUBTITLE = hex_to_rgb("2C3E50")
    COLOR_CREDITS = hex_to_rgb("556270")

    images_dir = config["paths"].get("images_dir", "output/images")
    mapping = config.get("background_mapping", {})
    overrides = config.get("layout_overrides", {})

    for slide_num, slide_raw in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # Background
        bg_path = resolve_bg_path(slide_num, mapping, images_dir)
        if os.path.exists(bg_path):
            slide.shapes.add_picture(os.path.abspath(bg_path), 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        else:
            print(f"Warning: Background missing for Slide {slide_num}: {bg_path}")
            
        if slide_num == 1:
            # Cover Slide (2 text boxes)
            cover_data = parse_cover_slide(slide_raw, lang_code)
            
            # TB 0: Title block
            tb_main = slide.shapes.add_textbox(Inches(0.223), Inches(1.098), Inches(5.038), Inches(3.545))
            tf_main = tb_main.text_frame
            tf_main.word_wrap = True
            tf_main.vertical_anchor = MSO_ANCHOR.TOP
            tf_main.margin_left = Inches(0.1)
            tf_main.margin_right = Inches(0.1)
            tf_main.margin_top = Inches(0.05)
            tf_main.margin_bottom = Inches(0.05)
            
            # Title
            p_title = tf_main.paragraphs[0]
            p_title.text = cover_data["title"]
            p_title.font.name = FONT_TITLE
            p_title.font.size = Pt(36 if lang_code == "ja" else 38)
            p_title.font.bold = True
            p_title.font.color.rgb = COLOR_TITLE
            p_title.space_after = Pt(8)
            p_title.alignment = PP_ALIGN.LEFT
            
            # Subtitle
            if cover_data["subtitle"]:
                p_sub = tf_main.add_paragraph()
                p_sub.text = cover_data["subtitle"]
                p_sub.font.name = FONT_TITLE
                p_sub.font.size = Pt(20 if lang_code == "ja" else 21)
                p_sub.font.italic = (lang_code != "ja")
                p_sub.font.color.rgb = COLOR_SUBTITLE
                p_sub.space_after = Pt(24)
                p_sub.alignment = PP_ALIGN.LEFT
            
            # Author
            p_author = tf_main.add_paragraph()
            p_author.text = cover_data["author"]
            p_author.font.name = FONT_AUTHOR
            p_author.font.size = Pt(14)
            p_author.font.bold = True
            p_author.font.color.rgb = COLOR_BODY
            p_author.space_after = Pt(18)
            p_author.alignment = PP_ALIGN.LEFT
            
            # TB 1: Credits block
            tb_cred = slide.shapes.add_textbox(Inches(0.223), Inches(6.234), Inches(6.797), Inches(1.001))
            tf_cred = tb_cred.text_frame
            tf_cred.word_wrap = True
            tf_cred.vertical_anchor = MSO_ANCHOR.TOP
            tf_cred.margin_left = Inches(0.1)
            tf_cred.margin_right = Inches(0.1)
            tf_cred.margin_top = Inches(0.05)
            tf_cred.margin_bottom = Inches(0.05)
            
            p_ai = tf_cred.paragraphs[0]
            p_ai.text = cover_data["ai_production"]
            p_ai.font.name = FONT_BODY
            p_ai.font.size = Pt(10)
            p_ai.font.italic = (lang_code != "ja")
            p_ai.font.color.rgb = COLOR_CREDITS
            p_ai.space_after = Pt(6)
            p_ai.alignment = PP_ALIGN.LEFT
            
            for cr_line in cover_data["credits"]:
                p_cr = tf_cred.add_paragraph()
                p_cr.text = cr_line
                p_cr.font.name = FONT_BODY
                p_cr.font.size = Pt(9.5)
                p_cr.font.color.rgb = COLOR_CREDITS
                p_cr.space_after = Pt(3)
                p_cr.alignment = PP_ALIGN.LEFT

        elif 2 <= slide_num <= config.get("content_slides_count", 50) + 1:
            # Content Slide (Single sentence)
            sentence = parse_content_slide(slide_raw)
            override = get_layout_override(slide_num, overrides)
            
            left = Inches(override["left_in"]) if override and "left_in" in override else Inches(0.90)
            top = Inches(override["top_in"]) if override and "top_in" in override else Inches(2.20)
            width = Inches(override["width_in"]) if override and "width_in" in override else Inches(4.70)
            height = Inches(override["height_in"]) if override and "height_in" in override else Inches(3.20)
            text_color = hex_to_rgb(override["text_color"]) if override and "text_color" in override else COLOR_BODY
            
            tb = slide.shapes.add_textbox(left, top, width, height)
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
            
            p = tf.paragraphs[0]
            p.text = sentence
            p.font.name = FONT_BODY
            p.font.bold = True
            
            char_len = len(sentence)
            if lang_code in ["en", "es"]:
                font_size = Pt(30 if char_len <= 55 else (28 if char_len <= 75 else 26))
            else:
                font_size = Pt(28 if char_len <= 22 else (26 if char_len <= 30 else 25))
                
            p.font.size = font_size
            p.font.color.rgb = text_color
            for r in p.runs:
                r.font.name = FONT_BODY
                r.font.size = font_size
                r.font.bold = True
                r.font.color.rgb = text_color
                
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.3

        else:
            # Vocabulary Slides
            title, words = parse_vocab_slide(slide_raw)
            
            # TB 0: Title
            tb_title = slide.shapes.add_textbox(Inches(0.850), Inches(0.750), Inches(5.800), Inches(0.421))
            tf_title = tb_title.text_frame
            tf_title.word_wrap = True
            tf_title.margin_left = Inches(0.1)
            tf_title.margin_right = Inches(0.1)
            tf_title.margin_top = Inches(0.05)
            tf_title.margin_bottom = Inches(0.05)
            
            p_title = tf_title.paragraphs[0]
            p_title.text = title
            p_title.font.name = FONT_BODY
            p_title.font.size = Pt(25)
            p_title.font.bold = True
            p_title.font.color.rgb = COLOR_TITLE
            p_title.space_after = Pt(14)
            p_title.alignment = PP_ALIGN.LEFT
            
            # TB 1: List
            list_w = Inches(7.056 if slide_num == 52 else 7.265)
            list_h = Inches(5.362 if slide_num == 52 else 5.826)
            tb_list = slide.shapes.add_textbox(Inches(0.850), Inches(1.320), list_w, list_h)
            tf_list = tb_list.text_frame
            tf_list.word_wrap = True
            tf_list.margin_left = Inches(0.1)
            tf_list.margin_right = Inches(0.1)
            tf_list.margin_top = Inches(0.05)
            tf_list.margin_bottom = Inches(0.05)
            
            for idx, item in enumerate(words):
                p_word = tf_list.paragraphs[0] if idx == 0 else tf_list.add_paragraph()
                p_word.space_after = Pt(5)
                p_word.line_spacing = 1.15
                
                parts = re.split(r"\s*[—–-]\s*", item, maxsplit=1)
                if len(parts) == 2:
                    headword, definition = parts[0], parts[1]
                    
                    r1 = p_word.add_run()
                    r1.text = headword + " — "
                    r1.font.name = FONT_BODY
                    r1.font.bold = True
                    r1.font.size = Pt(24 if lang_code != "ja" else 22)
                    r1.font.color.rgb = COLOR_TITLE
                    
                    r2 = p_word.add_run()
                    r2.text = definition
                    r2.font.name = FONT_BODY
                    r2.font.bold = False
                    r2.font.size = Pt(20 if slide_num == 52 and lang_code != "ja" else (22 if lang_code == "ja" else 24))
                    r2.font.color.rgb = COLOR_BODY
                else:
                    r = p_word.add_run()
                    r.text = item
                    r.font.name = FONT_BODY
                    r.font.size = Pt(24)
                    r.font.color.rgb = COLOR_BODY

    os.makedirs(os.path.dirname(output_pptx) or ".", exist_ok=True)
    prs.save(output_pptx)
    print(f"✓ Saved {output_pptx} ({len(prs.slides)} slides)")

def main():
    parser = argparse.ArgumentParser(description="Assemble PowerPoint presentations from text and image assets.")
    parser.add_argument("--config", default="templates/deck_config.template.json", help="Path to deck configuration JSON")
    args = parser.parse_args()

    with open(args.config, "r", encoding="utf-8") as f:
        config = json.load(f)

    paths = config["paths"]
    create_deck("en", paths["content_en"], paths["output_pptx_en"], config)
    create_deck("es", paths["content_es"], paths["output_pptx_es"], config)
    create_deck("ja", paths["content_ja"], paths["output_pptx_ja"], config)
    print("\nPresentation assembly complete!")

if __name__ == "__main__":
    main()
